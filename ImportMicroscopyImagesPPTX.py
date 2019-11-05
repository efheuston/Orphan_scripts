import os
import sys
import pptx
import re

# all images 3.02" high 2.99" wide
# lefts 0.28" 3.43" 6.58"
# tops  0.93" 4.24"

# one text box per image, always 0.2" higher
# height 0.25" width 1.5"
# Titles, clockwise from top left:
#  CD41 e450
#  Gr1 AF488
#  CD55 AF647
#  CellTrace Red/Orange
#  BrightField
#  Merge

# are there enough arguments
if len(sys.argv) != 3:
	print('Usage: ImportMicroscopyImagesPPTX.py <folder_to_scan> <output_powerpoint_name>')
	sys.exit()

# first argument is the folder to scan
scanFolder = sys.argv[1]

# make sure this is a valid path
if not os.path.isdir(scanFolder):
	print('The first argument does not appear to be a valid path. Exitting.')
	sys.exit()

# second argument is the path and name of the powerpoint output file
outputPptxName = sys.argv[2]

# make sure this name ends in pptx
if(outputPptxName[-5:].lower() != '.pptx'):
	outputPptxName = outputPptxName + '.pptx'

#debug output for the user
print('Reading images from the folder ' + scanFolder)
print('Writing results to ' + outputPptxName)

# clockwise from top-left, per image info
lefts = [0.28, 3.43, 6.58, 6.58, 3.43, 0.28]
tops = [0.93, 0.93, 0.93, 4.24, 4.24, 4.24]
titles = ['CD41 e450', 'Gr1 AF488', 'CD55 AF647', 'CellTrace Red/Orange', 'BrightField', 'Merge']
imageNameSuffix = ['_DAPI.tif', '_GFP.tif', '_CY5.tif', '_Rhodamine.tif', '_TL DIC.tif', '.tif']

# scan the input folder for czi files and use them to find the tif images
regex = re.compile('(.*czi$)')

cziList = []
for root, dirs, files in os.walk(scanFolder):
	for file in files:
		if regex.match(file):
			cziList.append(file[:-4]) # strip the .czi at the end

# debug print
#print(cziList)

# for each group, try to find the maching tifs
groupList = []
for root, dirs, files in os.walk(scanFolder):
	for czi in cziList:
		group = []
		for suffix in imageNameSuffix:
			regex = re.compile(czi + '-Image Export-\d+' + suffix)
			found = False
			for file in files:
				if regex.match(file):
					if not found:
						group.append(root + file)
						found = True
					else:
						print('Found multiple matching files for regex ' + regex.pattern)
						sys.exit()
			if not found:
				print('Did not find a match for regex ' + regex.pattern)
				sys.exit()
#			else:
#				print('Found a match for regex ' + regex.pattern) # debug
		groupList.append(group)

print('Found ' + str(len(groupList)) + ' sets of images')

# now start the process
# Create a new, empty, presentation
prs = pptx.Presentation()

# Apparently powerpoint always has the same slide layouts in the same order
# This one is the blank slide
blank_layout = prs.slide_layouts[6]

# image dimensions
imWidthInches = pptx.util.Inches(2.99)
imHeightInches = pptx.util.Inches(3.02)

# text dimensions
txTopOffsetInches = pptx.util.Inches(-0.2)
txWidthInches = pptx.util.Inches(1.25)
txHeightInches = pptx.util.Inches(0.25)

for group in groupList:
	# add a blank slide
	slide = prs.slides.add_slide(blank_layout)

	# make the slide background do its own thing
	#slide.follow_master_background = False

	# make the slide background solid black
	slide.background.fill.solid()
	slide.background.fill.fore_color.rgb = pptx.dml.color.RGBColor(0,0,0)
	
	for index in range(0,6):
		leftInches = pptx.util.Inches(lefts[index])
		topInches = pptx.util.Inches(tops[index])
		title = titles[index]
		
		# find the file which matches the target group and image
		img_path = group[index]
	
		# add an image to the slide
		pic = slide.shapes.add_picture(img_path, leftInches, topInches, imWidthInches, imHeightInches)
	
		# add a title box to the slide
		txBox = slide.shapes.add_textbox(leftInches, topInches+txTopOffsetInches, txWidthInches, txHeightInches)
		tf = txBox.text_frame
		tf.text = title
		#tf.font.color.rgb = pptx.dml.color.RGBColor(255,255,255)

# save the powerpoint file to disk
prs.save(outputPptxName)
